import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_file="AlgoSentinel_Deck.pptx"):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_COLOR = RGBColor(11, 15, 25)          # Deep dark slate #0B0F19
    CARD_BG = RGBColor(19, 27, 46)           # Card dark navy #131B2E
    CARD_BORDER = RGBColor(38, 52, 84)       # Card subtle border #263454
    CYAN = RGBColor(0, 229, 255)             # Accent Cyan #00E5FF
    EMERALD = RGBColor(0, 230, 118)          # Accent Green #00E676
    AMBER = RGBColor(255, 179, 0)            # Accent Amber #FFB300
    TEXT_WHITE = RGBColor(248, 250, 252)     # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)     # Slate 400 #94A3B8
    TEXT_DIM = RGBColor(100, 116, 139)       # Slate 500 #64748B
    PURPLE_ACCENT = RGBColor(168, 85, 247)   # AI accent

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, tag_color=CYAN):
        # Pill Tag / Badge
        tag_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.35))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = tag_color
        p.font.name = "Calibri"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.font.name = "Calibri"

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Center hero card
    card1 = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER
    card1.line.width = Pt(1.5)

    # Top Pill
    pill = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.416), Inches(1.8), Inches(4.5), Inches(0.45)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(16, 42, 60)
    pill.line.color.rgb = CYAN
    pill.line.width = Pt(1)
    ptf = pill.text_frame
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    pp.text = "ALPACA AI TRADING AGENTS HACKATHON"
    pp.font.size = Pt(11)
    pp.font.bold = True
    pp.font.color.rgb = CYAN
    pp.font.name = "Calibri"

    # Big Title
    title_box = slide1.shapes.add_textbox(Inches(1.8), Inches(2.55), Inches(9.733), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "AlgoSentinel 🛡️"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Calibri"

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.8), Inches(3.9), Inches(9.733), Inches(0.8))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Autonomous AI-Powered Options Trading Agent"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Calibri"

    # Team & Detail Info
    tag_box = slide1.shapes.add_textbox(Inches(1.8), Inches(4.8), Inches(9.733), Inches(1.0))
    tf = tag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Team AlgoSentinel\nPaper Trading Account  •  Real-Time News Sentiment  •  SPY Options Execution"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Calibri"

    # ==========================================
    # SLIDE 2: The Problem & Solution
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "🎯 Market Challenge & Core Innovation", "The Problem & Our Solution")

    # Left Card: The Problem
    card_prob = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(1.9), Inches(5.5), Inches(4.9)
    )
    card_prob.fill.solid()
    card_prob.fill.fore_color.rgb = CARD_BG
    card_prob.line.color.rgb = RGBColor(120, 40, 50)
    card_prob.line.width = Pt(1.5)

    tf_p = card_prob.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = Inches(0.4)
    tf_p.margin_top = Inches(0.35)
    tf_p.margin_right = Inches(0.4)

    p = tf_p.paragraphs[0]
    p.text = "⚠️  THE PROBLEM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 107, 107)

    p = tf_p.add_paragraph()
    p.text = "\nMost traders are emotional, biased, and slow."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf_p.add_paragraph()
    p.text = "\nOpportunities in market sentiment die in seconds."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(252, 165, 165)

    p = tf_p.add_paragraph()
    p.text = "\n• Human reaction latency misses fast-moving market catalyst headlines.\n• Cognitive bias & panic trading lead to broken discipline and blown accounts.\n• Manually screening options contracts, strikes, and ask prices takes too long."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # Right Card: The Solution
    card_sol = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.933), Inches(1.9), Inches(5.5), Inches(4.9)
    )
    card_sol.fill.solid()
    card_sol.fill.fore_color.rgb = CARD_BG
    card_sol.line.color.rgb = RGBColor(20, 100, 80)
    card_sol.line.width = Pt(1.5)

    tf_s = card_sol.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = Inches(0.4)
    tf_s.margin_top = Inches(0.35)
    tf_s.margin_right = Inches(0.4)

    p = tf_s.paragraphs[0]
    p.text = "⚡  OUR SOLUTION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    items = [
        ("AI reads live news in real-time", "Pulls S&P 500 & SPY RSS feeds every 15 minutes"),
        ("Scores sentiment automatically", "Gemini AI parses tone, nuances, and macro implications"),
        ("Executes trades instantly", "Direct Alpaca Trading API execution with optimal ATM strikes"),
        ("No human delays. No emotion.", "Pure signal-driven trading governed by strict risk gates")
    ]

    for title, desc in items:
        p = tf_s.add_paragraph()
        p.text = f"\n→ {title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p2 = tf_s.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: How It Works (Pipeline)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "🧠 Autonomous Trading Architecture", "How It Works: End-to-End Pipeline", CYAN)

    steps = [
        ("1. Google News", "Live SPY & S&P 500 macro headlines", CYAN),
        ("2. Gemini AI", "Scored BULLISH / BEARISH / NEUTRAL", PURPLE_ACCENT),
        ("3. Risk Gates", "5 safety checks (Floor, budget, pos limits)", AMBER),
        ("4. SPY Options", "ATM contract selection with live quotes", CYAN),
        ("5. Alpaca API", "Market order placement & live tracking", RGBColor(56, 189, 248)),
        ("6. ✅ Execution", "Position auto-managed (50% SL / 100% TP)", EMERALD)
    ]

    box_width = Inches(1.78)
    box_height = Inches(3.9)
    spacing = Inches(0.18)
    start_left = Inches(0.9)

    for i, (title, desc, accent) in enumerate(steps):
        x = start_left + i * (box_width + spacing)
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.9), box_width, box_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)

        # Step indicator on top
        step_badge = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.15), Inches(2.1), box_width - Inches(0.3), Inches(0.35)
        )
        step_badge.fill.solid()
        step_badge.fill.fore_color.rgb = RGBColor(16, 25, 45)
        step_badge.line.color.rgb = accent
        step_badge.line.width = Pt(1)
        sbtf = step_badge.text_frame
        sbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = sbtf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.text = f"STEP {i+1}"
        sp.font.size = Pt(10)
        sp.font.bold = True
        sp.font.color.rgb = accent

        # Step Text
        stf = card.text_frame
        stf.word_wrap = True
        stf.margin_top = Inches(0.7)
        stf.margin_left = Inches(0.12)
        stf.margin_right = Inches(0.12)

        p = stf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title.split(". ", 1)[1] if ". " in title else title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p2 = stf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # Bottom banner
    banner = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(6.1), Inches(11.533), Inches(0.7)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(18, 32, 54)
    banner.line.color.rgb = CYAN
    banner.line.width = Pt(1)
    btf = banner.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "⚡  Every 15 minutes. Fully autonomous. Zero human intervention needed."
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = CYAN

    # ==========================================
    # SLIDE 4: Why It Matters & Results
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "📊 Results, Rigor & Impact", "Why It Matters: Institutional Discipline for AI", EMERALD)

    col_width = Inches(3.64)
    col_gap = Inches(0.3)
    c_y = Inches(1.9)
    c_h = Inches(4.35)

    # Column 1: Autonomous & Transparent
    c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), c_y, col_width, c_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CARD_BORDER
    c1.line.width = Pt(1.5)

    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "🤖 AUTONOMOUS & AI-DRIVEN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN

    p = tf1.add_paragraph()
    p.text = "\n• Zero Human Latency"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf1.add_paragraph()
    p.text = "   Runs every 15 mins during market hours"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    p = tf1.add_paragraph()
    p.text = "\n• Gemini Sentiment Engine"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf1.add_paragraph()
    p.text = "   Gemini scores sentiment & drives decisions"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    p = tf1.add_paragraph()
    p.text = "\n• Complete Audit Trail"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf1.add_paragraph()
    p.text = "   Every trade and skip logged to trades.log"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # Column 2: 5 Risk Gates
    c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + col_width + col_gap), c_y, col_width, c_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = AMBER
    c2.line.width = Pt(1.5)

    tf2 = c2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = Inches(0.3)

    p = tf2.paragraphs[0]
    p.text = "🛡️ 5-TIER RISK MANAGEMENT"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AMBER

    gates = [
        ("Equity Floor", "$85k minimum capital floor"),
        ("Max 3 Positions", "Strict limit on simultaneous exposure"),
        ("5% Per-Trade Cap", "Max budget per individual contract"),
        ("Confidence Filter", "Only trade if signal confidence ≥ 60%"),
        ("Market Hours Only", "Active 9:30 AM – 4:00 PM ET")
    ]

    for name, val in gates:
        p = tf2.add_paragraph()
        p.text = f"\n✓ {name}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p2 = tf2.add_paragraph()
        p2.text = f"   {val}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # Column 3: Results & Stack
    c3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + (col_width + col_gap) * 2), c_y, col_width, c_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = CARD_BORDER
    c3.line.width = Pt(1.5)

    tf3 = c3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_right = tf3.margin_top = Inches(0.3)

    p = tf3.paragraphs[0]
    p.text = "📈 EXECUTION & RESULTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf3.add_paragraph()
    p.text = "\n• Fresh $100k Paper Account"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf3.add_paragraph()
    p.text = "   Live paper trading via Alpaca API"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    p = tf3.add_paragraph()
    p.text = "\n• Intelligent Position Exits"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf3.add_paragraph()
    p.text = "   -50% Stop-loss & +100% Take-profit protection"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    p = tf3.add_paragraph()
    p.text = "\n• Hackathon Stack"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p = tf3.add_paragraph()
    p.text = "   Alpaca Trading API • Google Gemini • Python"
    p.font.size = Pt(12)
    p.font.color.rgb = CYAN

    # Bottom pill
    foot = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(6.5), Inches(11.533), Inches(0.5)
    )
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(15, 30, 48)
    foot.line.color.rgb = CARD_BORDER
    foot.line.width = Pt(1)
    ftf = foot.text_frame
    ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fp.text = "Built on Alpaca Trading API + Google Gemini  •  Paper Trading: $100k Account  •  lablab.ai Hackathon"
    fp.font.size = Pt(11)
    fp.font.color.rgb = TEXT_MUTED

    prs.save(output_file)
    print(f"[SUCCESS] Deck saved to {output_file}")

if __name__ == "__main__":
    create_deck()
